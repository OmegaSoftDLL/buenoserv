#!/usr/bin/env python3
"""Gerador de apresentações PowerPoint (.pptx) — python-pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys, json

COR_AZUL = RGBColor(0x1A, 0x23, 0x7E)
COR_BRANCA = RGBColor(0xFF, 0xFF, 0xFF)

def nova_apresentacao(largura=13.333, altura=7.5):
    prs = Presentation()
    prs.slide_width = Inches(largura)
    prs.slide_height = Inches(altura)
    return prs

def salvar(prs, path):
    prs.save(path)
    print(f"PPTX|OK|{path}")

def slide_capa(prs, titulo, sub="", cliente="", data=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COR_AZUL
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(40)
    p.font.color.rgb = COR_BRANCA
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if sub:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = sub
        p2.font.size = Pt(24)
        p2.font.color.rgb = COR_BRANCA
        p2.alignment = PP_ALIGN.CENTER

def slide_conteudo(prs, titulo, itens):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COR_AZUL
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf2 = txBox2.text_frame
    for i, item in enumerate(itens):
        p2 = tf2.add_paragraph()
        p2.text = item
        p2.font.size = Pt(18)
        p2.space_after = Pt(8)

def slide_duas_colunas(prs, titulo, esquerda, direita):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COR_AZUL
    txL = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    tfL = txL.text_frame
    for item in esquerda:
        p2 = tfL.add_paragraph()
        p2.text = item
        p2.font.size = Pt(16)
    txR = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(6), Inches(5))
    tfR = txR.text_frame
    for item in direita:
        p2 = tfR.add_paragraph()
        p2.text = item
        p2.font.size = Pt(16)

def slide_tabela(prs, titulo, cabecalhos, dados):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COR_AZUL
    rows = len(dados) + 1
    cols = len(cabecalhos)
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5*rows)).table
    for i, h in enumerate(cabecalhos):
        table.cell(0, i).text = h
    for r, row_data in enumerate(dados, 1):
        for c, v in enumerate(row_data):
            table.cell(r, c).text = str(v)

def gerar_apresentacao(nome, titulo, slides_data, cliente="", data=""):
    prs = nova_apresentacao()
    slide_capa(prs, titulo, cliente=cliente, data=data)
    for s in slides_data:
        t = s.get('tipo', 'conteudo')
        if t == 'conteudo':
            slide_conteudo(prs, s['titulo'], s['itens'])
        elif t == 'duas_colunas':
            slide_duas_colunas(prs, s['titulo'], s['esquerda'], s['direita'])
        elif t == 'tabela':
            slide_tabela(prs, s['titulo'], s['cabecalhos'], s['dados'])
    salvar(prs, nome)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: gen_pptx.py <comando> <args_json>")
        print("Comandos: apresentacao")
        sys.exit(1)
    cmd = sys.argv[1]
    args = json.loads(sys.argv[2]) if len(sys.argv) > 2 else {}
    if cmd == "apresentacao":
        gerar_apresentacao(args["nome"], args["titulo"], args["slides"],
                         args.get("cliente",""), args.get("data",""))
    else:
        print(f"Comando desconhecido: {cmd}")
